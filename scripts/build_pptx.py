#!/usr/bin/env python3
"""Build the CafeFlow defense deck by reusing the prashant template in place.

It keeps the template's theme, banner, logos, fonts and bullet styling, and only
replaces the per-slide title text, body bullets and images with CafeFlow content.
"""
import copy
import sys
from pptx import Presentation
from pptx.oxml.ns import qn
from pptx.util import Emu

TEMPLATE = "/Users/nirajkafle/Downloads/prashant/report_docs/securelogti_presentation.pptx"
OUT = "/Users/nirajkafle/Downloads/amit/finalproject/docs/submission/presentation/CafeFlow_Presentation.pptx"
SHOTS = "/Users/nirajkafle/Downloads/amit/finalproject/docs/submission/assets/screenshots"
DIAG = "/Users/nirajkafle/Downloads/amit/finalproject/docs/submission/assets/diagrams"

# Old template titles, used to locate the title shape on each slide (index 0..14).
OLD_TITLES = [
    None,  # slide 0 = title slide (special)
    "Introduction",
    "Problem Statement",
    "Objectives",
    "Functional Requirements",
    "Non-functional Requirements",
    "Data Modeling: ER Diagram",
    "Process Modeling: Data Flow Diagrams",
    "System Architecture",
    "Algorithm Details",
    "Implementing Tools",
    "Implementation: Module Details",
    "Application Demo: Dashboard",
    "Application Demo: Threat Intelligence",
    "Conclusion",
]

# New title + body bullets (text, level) per slide. None body = image slide.
NEW = {
    1: ("Introduction", [
        ("Cafes handle many time-critical daily tasks: taking orders, tracking them, billing, and managing stock.", 0),
        ("Most small cafes still do this manually with notebooks, verbal orders, and spreadsheets.", 0),
        ("This is slow, error-prone, and gives no reliable view of sales or stock.", 1),
        ("CafeFlow: a web-based system that unifies ordering, billing, inventory, staff and reporting.", 0),
    ]),
    2: ("Problem Statement", [
        ("Manual cafe operations cause real, everyday problems:", 0),
        ("Paper orders relayed verbally lead to delays and mistakes.", 1),
        ("Bills calculated by hand are error-prone and cause disputes.", 1),
        ("Stock tracked informally runs out unexpectedly.", 1),
        ("No usable sales data, so owners lack insight into revenue and best-sellers.", 1),
        ("Commercial POS systems are costly, complex, or hardware-bound.", 1),
        ("Need: a simple, affordable, web-based system built for cafes.", 0),
    ]),
    3: ("Objectives", [
        ("General: a single web platform to digitise cafe ordering, billing, inventory and staff management.", 0),
        ("Specific objectives:", 0),
        ("Provide secure, role-based access for Admin and Staff.", 1),
        ("Manage menu, orders, billing and inventory in real time.", 1),
        ("Track orders through a pending → preparing → completed workflow.", 1),
        ("Generate accurate sales and inventory analytics.", 1),
    ]),
    4: ("Functional Requirements", [
        ("Email/password login with role-based access (Admin vs Staff).", 0),
        ("Menu management: create, edit, delete, toggle availability.", 0),
        ("Order taking with table/customer details and live status tracking.", 0),
        ("Itemised billing for orders, ready to print.", 0),
        ("Inventory tracking with automatic low-stock alerts.", 0),
        ("Staff management and sales reporting from real order data.", 0),
    ]),
    5: ("Non-functional Requirements", [
        ("Security: bcrypt-hashed passwords and HTTP-only session cookies.", 0),
        ("Usability: clean, responsive UI on desktop, tablet and mobile.", 0),
        ("Performance: common actions complete in a fraction of a second.", 0),
        ("Reliability: data persisted in SQLite, surviving restarts.", 0),
        ("Maintainability: clear layers — UI, API and data access.", 0),
    ]),
    6: ("Data Modeling: ER Diagram", f"{DIAG}/er.png"),
    7: ("Process Modeling: Order Flow", f"{DIAG}/workflow.png"),
    8: ("System Architecture", f"{DIAG}/architecture.png"),
    9: ("Key Logic: Orders & Reports", [
        ("Order workflow: each order moves pending → preparing → completed; each line stores a price snapshot.", 0),
        ("Billing: totals computed from order lines; an itemised bill is generated for printing.", 0),
        ("Inventory: items at or below their threshold are flagged as low stock.", 0),
        ("Reports: revenue, orders, top items and category split are aggregated from stored orders.", 0),
        ("Auth: passwords hashed with bcrypt; every API request validates the session and role.", 0),
    ]),
    10: ("Implementing Tools", [
        ("Next.js + React + TypeScript: full-stack UI and API routes.", 0),
        ("SQLite via better-sqlite3: embedded relational database.", 0),
        ("Tailwind CSS + shadcn/ui: styling and accessible components.", 0),
        ("Zustand: lightweight client-side state management.", 0),
        ("bcrypt: password hashing; Recharts: dashboard charts.", 0),
        ("Playwright: end-to-end testing; Git/GitHub + VS Code.", 0),
    ]),
    11: ("Implementation: Module Details", [
        ("Authentication: login, hashed passwords, sessions, role-based access.", 0),
        ("Menu: CRUD, availability toggle, search, filter, table/card views.", 0),
        ("Orders: build cart, submit, track status; price snapshot per line.", 0),
        ("Billing: itemised bills for completed and in-progress orders.", 0),
        ("Inventory: stock with units and low-stock thresholds + alerts.", 0),
        ("Staff & Reports: team management and analytics from real orders.", 0),
    ]),
    12: ("Application Demo: Dashboard", f"{SHOTS}/02-dashboard.png"),
    13: ("Application Demo: Orders", f"{SHOTS}/04-orders-new.png"),
    14: ("Conclusion", [
        ("CafeFlow is a complete, working cafe management system: menu, orders, billing, inventory, staff and reports.", 0),
        ("Replaces manual processes with a fast, consistent, role-based digital workflow.", 0),
        ("Secure authentication, SQLite persistence and analytics from real order data.", 0),
        ("Validated by 25 automated end-to-end tests, all passing.", 0),
        ("Future: online payments, multi-outlet, kitchen display and QR table ordering.", 0),
    ]),
}

# Title slide (index 0) paragraph lines, mapped onto existing paragraphs in order.
TITLE_LINES = [
    "CafeFlow",
    "A Web-Based Cafe Management System",
    "for Small and Medium Cafes",
    "Amit Tharu, BIT (LC0003001642)",
    "Supervisor: Mr. Saishab Bhattarai",
    "Phoenix College of Management",
]


def set_runs_text(p_el, text):
    rs = p_el.findall(qn('a:r'))
    if not rs:
        return False
    t = rs[0].find(qn('a:t'))
    if t is None:
        t = rs[0].makeelement(qn('a:t'), {})
        rs[0].append(t)
    t.text = text
    for r in rs[1:]:
        p_el.remove(r)
    return True


def para_level(p_el):
    pPr = p_el.find(qn('a:pPr'))
    if pPr is not None and pPr.get('lvl'):
        return int(pPr.get('lvl'))
    return 0


def rebuild_body(tf, bullets):
    txBody = tf._txBody
    paras = txBody.findall(qn('a:p'))
    templ = {}
    for p in paras:
        lvl = para_level(p)
        templ.setdefault(lvl, copy.deepcopy(p))
    default = templ.get(0, copy.deepcopy(paras[0]))
    for p in paras:
        txBody.remove(p)
    markers = {0: "•  ", 1: "◦  "}  # bullets are literal text in this template
    for text, lvl in bullets:
        src = templ.get(lvl, default)
        newp = copy.deepcopy(src)
        set_runs_text(newp, markers.get(lvl, "•  ") + text)
        txBody.append(newp)


def text_shapes(slide):
    return [s for s in slide.shapes if s.has_text_frame and s.text_frame.text.strip()]


def picture_shapes(slide):
    return [s for s in slide.shapes if s.shape_type == 13 or s._element.tag.endswith('}pic')]


def find_title(slide, old_title):
    for s in text_shapes(slide):
        if s.text_frame.text.strip() == old_title:
            return s
    # fallback: topmost text shape
    ts = sorted(text_shapes(slide), key=lambda s: s.top or 0)
    return ts[0] if ts else None


def replace_picture(slide, new_path):
    pics = picture_shapes(slide)
    if not pics:
        return
    first = pics[0]
    left, top, width = first.left, first.top, first.width
    for p in pics:  # remove all existing pictures
        p._element.getparent().remove(p._element)
    slide.shapes.add_picture(new_path, left, top, width=width)


prs = Presentation(TEMPLATE)
slides = list(prs.slides)
print(f"template slides: {len(slides)}")

# ---- Slide 0 : title slide ----
s0 = slides[0]
box = max(text_shapes(s0), key=lambda s: len(s.text_frame.text))
paras = box.text_frame._txBody.findall(qn('a:p'))
for i, p in enumerate(paras):
    set_runs_text(p, TITLE_LINES[i] if i < len(TITLE_LINES) else "")
print("slide 1: title set")

# ---- Slides 1..14 ----
for idx in range(1, len(slides)):
    if idx not in NEW:
        continue
    slide = slides[idx]
    title, body = NEW[idx]
    t = find_title(slide, OLD_TITLES[idx])
    if t is not None:
        set_runs_text(t.text_frame._txBody.find(qn('a:p')), title)
    if isinstance(body, str):  # image slide
        # blank any non-title text frames (e.g. "Level 0 / Level 1" labels)
        for s in text_shapes(slide):
            if t is None or s._element is not t._element:
                for p in s.text_frame._txBody.findall(qn('a:p')):
                    set_runs_text(p, "")
        replace_picture(slide, body)
    else:  # bullet slide
        bodies = [s for s in text_shapes(slide) if t is None or s._element is not t._element]
        if bodies:
            target = max(bodies, key=lambda s: len(s.text_frame.text))
            rebuild_body(target.text_frame, body)
    print(f"slide {idx+1}: {title}")

prs.save(OUT)
print("saved ->", OUT)
